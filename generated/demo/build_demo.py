# -*- coding: utf-8 -*-
"""Демо-комплект для видео про CueDeck.
Презентации сами ведут рассказ: переключая файлы, Азат идёт по своей теме.
Форматы намеренно разные — PDF, PPTX, картинки, видео."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = "/Users/azatlife/VSCode/pdf-presenter/generated/demo"
os.makedirs(OUT, exist_ok=True)

NAVY   = RGBColor(0x14, 0x23, 0x3A)
ACCENT = RGBColor(0xE8, 0x6A, 0x17)
GRAY   = RGBColor(0x6B, 0x74, 0x80)
LIGHT  = RGBColor(0xF4, 0xF5, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xA8, 0xB4, 0xC4)
DIM    = RGBColor(0x5A, 0x67, 0x78)
FONT   = "Arial"
BRAND  = "CueDeck"

W, H = Inches(13.333), Inches(7.5)


def new_prs():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = bg
    return s


def no_shadow(sh):
    """LibreOffice подрисовывает тень по умолчанию — гасим явным пустым effectLst."""
    sh.shadow.inherit = False
    el = sh._element
    spPr = el.spPr
    for e in spPr.findall(qn('a:effectLst')):
        spPr.remove(e)
    spPr.append(spPr.makeelement(qn('a:effectLst'), {}))
    for st in el.findall(qn('p:style')):
        el.remove(st)
    return sh


def tb(slide, l, t, w, h, text, size=24, color=NAVY, bold=False,
       align=PP_ALIGN.LEFT, spacing=1.15, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name, f.size, f.bold, f.color.rgb = FONT, Pt(size), bold, color
    return box


def rect(slide, l, t, w, h, color=ACCENT, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return no_shadow(sh)


def outline(slide, l, t, w, h, color=ACCENT, width=2.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.background()
    sh.line.color.rgb = color
    sh.line.width = Pt(width)
    sh.adjustments[0] = 0.04
    return no_shadow(sh)


def badge(slide, text, dark=False):
    """Плашка формата в углу — на камере сразу видно, что файлы разные."""
    tb(slide, 11.05, 0.55, 1.6, 0.4, text, 11,
       ACCENT if dark else GRAY, bold=True, align=PP_ALIGN.RIGHT)


def footer(slide, n, total, dark=False):
    c = DIM if dark else GRAY
    tb(slide, 0.7, 6.85, 7, 0.4, BRAND, 10, c)
    tb(slide, 11.0, 6.85, 1.6, 0.4, f"{n} / {total}", 10, c, align=PP_ALIGN.RIGHT)


# ---------- типовые слайды ----------

def s_title(prs, title, subtitle, fmt):
    s = blank(prs, NAVY)
    rect(s, 0.7, 1.5, 1.6, 0.09)
    tb(s, 0.7, 0.55, 8, 0.4, BRAND, 12, MUTED)
    badge(s, fmt, dark=True)
    tb(s, 0.7, 1.9, 11.5, 2.0, title, 46, WHITE, bold=True, spacing=1.1)
    tb(s, 0.7, 4.0, 10.5, 0.8, subtitle, 20, MUTED, spacing=1.3)
    rect(s, 0.7, 5.3, 0.06, 0.9, ACCENT)
    tb(s, 1.0, 5.3, 8, 0.5, "Азат Хусаенов", 20, WHITE, bold=True)
    tb(s, 1.0, 5.77, 8, 0.5, "Технический директор мероприятий", 14, MUTED)
    return s


def s_section(prs, kicker, title, n, total):
    s = blank(prs, LIGHT)
    rect(s, 0.7, 2.9, 1.2, 0.07)
    tb(s, 0.7, 2.35, 8, 0.4, kicker, 13, ACCENT, bold=True)
    tb(s, 0.7, 3.2, 11, 1.6, title, 40, NAVY, bold=True, spacing=1.1)
    footer(s, n, total)
    return s


def s_bullets(prs, title, items, n, total):
    s = blank(prs)
    tb(s, 0.7, 0.75, 11.5, 1.0, title, 32, NAVY, bold=True)
    rect(s, 0.7, 1.72, 1.0, 0.06)
    top = 2.3
    for it in items:
        rect(s, 0.75, top + 0.17, 0.12, 0.12, ACCENT)
        tb(s, 1.15, top, 10.8, 0.9, it, 18, RGBColor(0x33, 0x3B, 0x47), spacing=1.25)
        top += 0.88
    footer(s, n, total)
    return s


def s_big(prs, number, caption, note, n, total):
    s = blank(prs, NAVY)
    tb(s, 0.7, 2.0, 11.9, 1.8, number, 96, WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, 0.7, 3.9, 11.9, 0.7, caption, 26, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    tb(s, 2.2, 4.7, 8.9, 1.1, note, 15, MUTED, align=PP_ALIGN.CENTER, spacing=1.35)
    footer(s, n, total, dark=True)
    return s


def s_quote(prs, text, author, n, total):
    s = blank(prs, LIGHT)
    tb(s, 0.75, 1.15, 1.6, 1.3, "«", 72, RGBColor(0xD8, 0xDD, 0xE4), bold=True)
    tb(s, 1.95, 2.15, 10.2, 2.6, text, 27, NAVY, spacing=1.35)
    rect(s, 1.95, 5.05, 0.06, 0.7, ACCENT)
    tb(s, 2.25, 5.05, 9, 0.7, author, 15, GRAY, spacing=1.25)
    footer(s, n, total)
    return s


def s_three_windows(prs, n, total):
    """Схема трёх окон — главный слайд всего рассказа."""
    s = blank(prs)
    tb(s, 0.7, 0.75, 11.5, 1.0, "Три окна вместо одного экрана", 32, NAVY, bold=True)
    rect(s, 0.7, 1.72, 1.0, 0.06)
    cols = [
        ("ОПЕРАТОР", "Текущий слайд\nи следующий\n\nПлейлист спикеров\nТаймеры и заметки", ACCENT),
        ("СУФЛЁР", "Слайд спикера\nи что дальше\n\nТаймер выступления\nСообщения от меня", NAVY),
        ("ЗАЛ", "Только слайд\n\nБез курсора,\nбез панелей,\nбез alt-tab", GRAY),
    ]
    x = 0.75
    for name, body, color in cols:
        outline(s, x, 2.3, 3.75, 3.9, color, 2.0)
        rect(s, x, 2.3, 3.75, 0.5, color)
        tb(s, x, 2.36, 3.75, 0.4, name, 14, WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + 0.35, 3.05, 3.05, 3.0, body, 15, RGBColor(0x33, 0x3B, 0x47), spacing=1.3)
        x += 4.05
    footer(s, n, total)
    return s


def s_embedded_video(prs, n, total):
    """Видео, вшитое прямо в слайд PPTX — CueDeck достаёт его и играет оверлеем."""
    s = blank(prs)
    tb(s, 0.7, 0.6, 11.5, 0.9, "Видео прямо в слайде", 32, NAVY, bold=True)
    rect(s, 0.7, 1.55, 1.0, 0.06)
    tb(s, 0.7, 1.85, 11.5, 0.5,
       "Ролик лежит внутри этого PPTX. Программа достаёт его и играет поверх слайда.",
       16, GRAY)
    movie = f"{OUT}/04_Ролик-о-программе.mp4"
    poster = f"{OUT}/_постер.png"
    if os.path.exists(movie) and os.path.exists(poster):
        s.shapes.add_movie(movie, Inches(2.85), Inches(2.5), Inches(7.6), Inches(4.28),
                           poster_frame_image=poster, mime_type="video/mp4")
    else:
        outline(s, 2.85, 2.5, 7.6, 4.28, GRAY, 1.5)
        tb(s, 2.85, 4.4, 7.6, 0.5, "видеофайл не найден при сборке", 14, GRAY,
           align=PP_ALIGN.CENTER)
    footer(s, n, total)
    return s


def s_end(prs, title, lines, n, total):
    s = blank(prs, NAVY)
    tb(s, 0.7, 2.5, 11.5, 1.2, title, 42, WHITE, bold=True)
    rect(s, 0.7, 3.9, 1.2, 0.07)
    tb(s, 0.7, 4.35, 11, 1.6, lines, 18, MUTED, spacing=1.45)
    footer(s, n, total, dark=True)
    return s


# ---------- 01 · PDF · Проблема ----------
def deck_problem():
    p, T = new_prs(), 8
    s_title(p, "Почему я написал\nсвою программу",
            "Показ презентаций на живых мероприятиях", "PDF")
    s_bullets(p, "Обычный день на площадке", [
        "Восемь спикеров, у каждого свой файл",
        "PDF, PowerPoint, Keynote и картинки вперемешку",
        "Кто-то приносит флешку за две минуты до выхода",
        "Правки прилетают прямо во время выступления",
    ], 2, T)
    s_bullets(p, "Чего не хватало в обычной смотрелке", [
        "Не видно, сколько слайдов осталось и какой следующий",
        "Спикеру на сцене нечего показать: ни таймера, ни заметок",
        "Alt-tab между файлами видит весь зал",
        "Между выступлениями нечем закрыть экран",
    ], 3, T)
    s_quote(p, "Вот это вот постоянно такой вот гемор непонятный,\nкоторый происходит. Я решил попробовать это исправить.",
            "Из первого обзора CueDeck", 4, T)
    s_section(p, "ЗАДАЧА", "Собрать всё в одно окно\nи ничего лишнего не показать залу", 5, T)
    s_bullets(p, "Что должно быть у оператора", [
        "Все презентации в одной очереди, без переключения программ",
        "Видно текущий и следующий слайд",
        "Управление таймером и подсказками спикеру",
        "Экран зала под полным контролем",
    ], 6, T)
    s_big(p, "3 окна", "оператор · суфлёр · зал",
          "Дальше показываю, как это устроено внутри.", 7, T)
    s_end(p, "Следующий файл — PPTX",
          "Специально держу презентации в разных форматах.\nВ плейлисте они лежат вперемешку и открываются одинаково.", 8, T)
    p.save(f"{OUT}/01_Проблема.pptx")


# ---------- 02 · PPTX · Три окна ----------
def deck_windows():
    p, T = new_prs(), 8
    s_title(p, "Три окна", "Как устроен показ изнутри", "PPTX")
    s_three_windows(p, 2, T)
    s_embedded_video(p, 3, T)
    s_bullets(p, "Окно оператора — то, что вижу я", [
        "Текущий слайд крупно, следующий — рядом",
        "Сколько слайдов осталось до конца выступления",
        "Плейлист: кто выступает сейчас и кто за ним",
        "Заметки к слайду, которые в зал не уходят",
    ], 4, T)
    s_bullets(p, "Суфлёр — то, что видит спикер", [
        "Свой слайд и следующий, чтобы не оборачиваться на экран",
        "Таймер выступления: сколько осталось",
        "Сообщения от меня: «две минуты», «закругляйся»",
        "Заметки к выступлению, если спикер их прислал",
    ], 5, T)
    s_bullets(p, "Зал — то, что видят зрители", [
        "Только слайд, во весь экран",
        "Ни курсора, ни панелей, ни рабочего стола",
        "Заставка вместо чёрного экрана между спикерами",
        "Переключение без мигания и без чужих окон",
    ], 6, T)
    s_quote(p, "Смысл простой: на экране зала не должно появиться\nничего, чего я туда сознательно не отправил.",
            "Главное правило программы", 7, T)
    s_end(p, "Дальше — картинки",
          "Следующие файлы в плейлисте — обычные JPG.\nОткрываются так же, как презентации.", 8, T)
    p.save(f"{OUT}/02_Три-окна.pptx")


# ---------- 03 · JPG · Что ещё умеет ----------
def deck_features():
    p, T = new_prs(), 5
    s_title(p, "Что ещё умеет", "Мелочи, из-за которых всё и затевалось", "JPG")
    s_bullets(p, "Плейлист спикеров", [
        "Файлы кидаются перетаскиванием, порядок меняется мышкой",
        "Флешка за две минуты до выхода — просто ещё одна строка",
        "Файл переехал или пропал — программа скажет, а не упадёт",
    ], 2, T)
    s_bullets(p, "Таймеры и связь со спикером", [
        "Пресеты таймера по правой кнопке",
        "Шесть готовых сообщений на суфлёр в одно нажатие",
        "Таймер видит спикер, зал его не видит",
    ], 3, T)
    s_bullets(p, "Живой вход и заставка", [
        "Картинка с HDMI-капчера как обычный элемент плейлиста",
        "Заставка между выступлениями одной кнопкой",
        "Предпрослушка звука в наушники до вывода в зал",
    ], 4, T)
    s_end(p, "Дальше — видео",
          "Видеофайлы лежат в том же плейлисте.\nMP4, MOV, M4V, WEBM.", 5, T)
    p.save(f"{OUT}/03_Что-умеет.pptx")


# ---------- Заставка ----------
def splash():
    p = new_prs()
    s = blank(p, NAVY)
    rect(s, 5.9, 2.55, 1.5, 0.09)
    tb(s, 0.7, 2.9, 11.9, 1.3, BRAND, 60, WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, 0.7, 4.2, 11.9, 0.7, "Следующее выступление начнётся через несколько минут",
       20, MUTED, align=PP_ALIGN.CENTER)
    p.save(f"{OUT}/_заставка.pptx")


# ---------- Карточки для видеоролика ----------
def video_cards():
    p = new_prs()
    cards = [
        ("Восемь спикеров", "Восемь файлов. Три формата."),
        ("Один экран зала", "И ничего лишнего на нём."),
        ("Три окна", "Оператор. Суфлёр. Зал."),
        ("Таймер спикеру", "И шесть сообщений в одно нажатие."),
        ("Бесплатно", "Mac и Windows. Исходники открыты."),
        (BRAND, "azat.life"),
    ]
    for i, (title, sub) in enumerate(cards):
        s = blank(p, NAVY)
        rect(s, 5.9, 2.75, 1.5, 0.08)
        tb(s, 0.7, 3.1, 11.9, 1.1, title, 54, WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, 0.7, 4.35, 11.9, 0.7, sub, 22,
           ACCENT if i == len(cards) - 1 else MUTED, align=PP_ALIGN.CENTER)
    p.save(f"{OUT}/_карточки-ролика.pptx")


deck_problem(); deck_windows(); deck_features(); splash(); video_cards()
print("PPTX пересобраны")
